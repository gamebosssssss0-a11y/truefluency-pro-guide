"""
StudySprint Backend — FastAPI
Render server that connects the TrueFluency Pro frontend to Gemini Flash.

Flow:
  Frontend uploads file → Supabase stores it + calls extract edge fn
  Frontend calls /generate-mock → this backend fetches extracted text
  from Supabase → sends to Gemini → returns real questions to frontend
"""

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
import os
import json
import io
import httpx
from dotenv import load_dotenv

# File extraction libraries
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

load_dotenv()

app = FastAPI(title="StudySprint Backend", version="1.0.0")

# ── CORS ──────────────────────────────────────────────────────────────────────
# This allows the React frontend to call this backend.
# In production, replace "*" with your actual Lovable/frontend URL.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],        # TODO: lock down to your frontend URL in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── ENV VARS ──────────────────────────────────────────────────────────────────
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_KEY")   # service role key, NOT the publishable one
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OPENROUTER_MODEL = "nvidia/nemotron-3-ultra-550b-a55b:free"                           # fast + cheap for test phase


# ── REQUEST / RESPONSE MODELS ─────────────────────────────────────────────────

class MockRequest(BaseModel):
    material_id: str            # ID from course_materials table in Supabase
    course_code: str            # e.g. "GST111"
    course_name: str            # e.g. "Use of English"
    question_count: int = 20
    difficulty: str = "balanced"  # gentle | balanced | challenging | exam
    topic_focus: list[str] = []   # optional list of topics to weight heavily

    # User context — built from their onboarding profile
    user_goal: Optional[str] = None         # "pass" | "top-grades" | "catch-up"
    user_timeline: Optional[str] = None     # "lt-week" | "2-4-weeks" | "gt-month"
    user_level: Optional[str] = None        # "100" | "200" | "300" | "400" | "500"
    user_department: Optional[str] = None


class PredictRequest(BaseModel):
    material_id: str
    course_code: str
    course_name: str
    user_level: Optional[str] = None
    user_department: Optional[str] = None


# ── HELPERS ───────────────────────────────────────────────────────────────────

async def fetch_extracted_text(material_id: str) -> str:
    """
    Fetch the extracted_content from Supabase for a given material_id.
    Uses the service role key so it can bypass RLS.
    """
    headers = {
        "apikey": SUPABASE_SERVICE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_KEY}",
        "Content-Type": "application/json",
    }
    url = f"{SUPABASE_URL}/rest/v1/course_materials?id=eq.{material_id}&select=extracted_content,file_name,course_code,extraction_status"

    async with httpx.AsyncClient() as client:
        res = await client.get(url, headers=headers)

    if res.status_code != 200:
        raise HTTPException(status_code=502, detail=f"Supabase error: {res.text}")

    rows = res.json()
    if not rows:
        raise HTTPException(status_code=404, detail="Material not found")

    row = rows[0]
    if row.get("extraction_status") not in ("success",):
        raise HTTPException(
            status_code=422,
            detail=f"Material text not ready yet. Status: {row.get('extraction_status')}"
        )

    content = row.get("extracted_content", "")
    if not content or len(content.strip()) < 50:
        raise HTTPException(status_code=422, detail="Extracted content is too short to generate questions from.")

    return content


async def call_gemini(prompt: str) -> str:
    """
    Call Gemini Flash via OpenRouter and return the raw text response.
    """
    url = "https://openrouter.ai/api/v1/chat/completions"

    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
        "HTTP-Referer": "http://localhost:8000",
        "X-Title": "StudySprint",
    }

    body = {
        "model": OPENROUTER_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.7,
        "max_tokens": 4096,
    }

    async with httpx.AsyncClient(timeout=60.0) as client:
        res = await client.post(url, headers=headers, json=body)

    if res.status_code != 200:
        raise HTTPException(status_code=502, detail=f"Gemini error: {res.text}")

    data = res.json()
    try:
        return data["choices"][0]["message"]["content"]
    except (KeyError, IndexError):
        raise HTTPException(status_code=502, detail="Unexpected OpenRouter response format")

def build_user_context(req) -> str:
    """
    Build a user context string from the profile data sent by the frontend.
    This is what makes the AI understand WHO it's generating for.
    """
    goal_map = {
        "pass": "just needs to pass (focus on high-frequency topics only)",
        "top-grades": "aiming for top grades (cover everything including edge topics)",
        "catch-up": "catching up on missed content (cover fundamentals thoroughly)",
    }
    timeline_map = {
        "lt-week": "less than a week until their exam (urgent — prioritize the most likely topics)",
        "2-4-weeks": "2 to 4 weeks until their exam",
        "gt-month": "more than a month until their exam (can afford more breadth)",
        "unsure": "unsure of their exam date",
    }

    parts = []
    if req.user_level:
        parts.append(f"- Academic level: {req.user_level}L at University of Ibadan")
    if req.user_department:
        parts.append(f"- Department: {req.user_department}")
    if req.user_goal:
        parts.append(f"- Goal: Student {goal_map.get(req.user_goal, req.user_goal)}")
    if req.user_timeline:
        parts.append(f"- Timeline: {timeline_map.get(req.user_timeline, req.user_timeline)}")

    if not parts:
        return ""

    return "Student context:\n" + "\n".join(parts)


def difficulty_instruction(difficulty: str) -> str:
    return {
        "gentle": "Make the questions straightforward and accessible. Focus on definitions, basic concepts, and direct recall.",
        "balanced": "Mix easy recall questions with questions that require understanding and application.",
        "challenging": "Include questions that require deeper understanding, analysis, and application of concepts.",
        "exam": "Match real exam difficulty exactly. Include tricky options, edge cases, and questions requiring synthesis.",
    }.get(difficulty, "Mix easy and hard questions.")


# ── ROUTES ────────────────────────────────────────────────────────────────────

async def fetch_file_from_supabase(file_path: str) -> bytes:
    """Download a file from Supabase storage and return raw bytes."""
    url = f"{SUPABASE_URL}/storage/v1/object/course-materials/{file_path}"
    headers = {
        "apikey": SUPABASE_SERVICE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_KEY}",
    }
    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(url, headers=headers)
    if res.status_code != 200:
        raise HTTPException(status_code=502, detail=f"Could not fetch file from storage: {res.text}")
    return res.content


def extract_text_from_pdf(file_bytes: bytes) -> str:
    if not HAS_PDF:
        raise HTTPException(status_code=500, detail="PyPDF2 not installed")
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        text += (page.extract_text() or "") + "\n"
    return text.strip()


def extract_text_from_docx(file_bytes: bytes) -> str:
    if not HAS_DOCX:
        raise HTTPException(status_code=500, detail="python-docx not installed")
    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    return text.strip()


def extract_text_from_pptx(file_bytes: bytes) -> str:
    if not HAS_PPTX:
        raise HTTPException(status_code=500, detail="python-pptx not installed")
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text.strip()


async def save_extracted_content(material_id: str, content: str, status: str = "success"):
    """Save extracted text back to Supabase course_materials table."""
    headers = {
        "apikey": SUPABASE_SERVICE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_KEY}",
        "Content-Type": "application/json",
    }
    url = f"{SUPABASE_URL}/rest/v1/course_materials?id=eq.{material_id}"
    body = {"extracted_content": content, "extraction_status": status}
    async with httpx.AsyncClient() as client:
        res = await client.patch(url, headers=headers, json=body)
    if res.status_code not in (200, 204):
        raise HTTPException(status_code=502, detail=f"Failed to save extraction: {res.text}")


@app.get("/health")
def health():
    """Quick check that the server is alive — call this first when testing."""
    return {"status": "ok", "service": "StudySprint Backend"}


@app.post("/predict-topics")
async def predict_topics(req: PredictRequest):
    """
    Step 1 of the AI pipeline.
    Takes a material_id, fetches the extracted text from Supabase,
    sends it to Gemini, returns predicted exam topics with confidence scores.

    Frontend uses this to replace the hardcoded placeholder topics on the
    dashboard and course detail screens.
    """
    extracted_text = await fetch_extracted_text(req.material_id)

    # Trim to avoid token overload — first 6000 chars is usually enough for topic prediction
    trimmed = extracted_text[:6000]

    prompt = f"""
You are an exam prediction engine for University of Ibadan ({req.user_level or "undergraduate"} level).

Analyze this course material for {req.course_code} — {req.course_name}:
{req.user_department and f"Department: {req.user_department}" or ""}

--- MATERIAL START ---
{trimmed}
--- MATERIAL END ---

Based on this content, identify the most likely exam topics.
For each topic, give a confidence score (0.0 to 1.0) based on how prominently it features in the material.

Respond ONLY with a JSON array. No explanation, no markdown, no backticks. Just raw JSON like:
[
  {{"topic": "Topic Name", "confidence": 0.9}},
  {{"topic": "Another Topic", "confidence": 0.6}}
]

Return 5 to 8 topics, ordered from highest to lowest confidence.
"""

    raw = await call_gemini(prompt)

    # Clean up any accidental markdown the model adds
    cleaned = raw.strip().replace("```json", "").replace("```", "").strip()

    try:
        topics = json.loads(cleaned)
    except json.JSONDecodeError:
        raise HTTPException(status_code=502, detail=f"Gemini returned invalid JSON: {cleaned[:300]}")

    return {
        "course_code": req.course_code,
        "topics": topics,
        "material_id": req.material_id,
    }


@app.post("/generate-mock")
async def generate_mock(req: MockRequest):
    """
    Step 2 of the AI pipeline — the core endpoint.
    Generates a full mock test based on the extracted material.

    This replaces the hardcoded sampleQuestions from questions.ts in the frontend.
    """
    extracted_text = await fetch_extracted_text(req.material_id)

    # Trim to keep within Gemini's context window comfortably
    trimmed = extracted_text[:8000]

    user_ctx = build_user_context(req)
    diff_instruction = difficulty_instruction(req.difficulty)

    topic_instruction = ""
    if req.topic_focus:
        topic_instruction = f"Pay extra attention to these topics (weight them more heavily): {', '.join(req.topic_focus)}."

    prompt = f"""
You are a mock exam generator for University of Ibadan.

{user_ctx}

Course: {req.course_code} — {req.course_name}
Difficulty: {req.difficulty}. {diff_instruction}
{topic_instruction}

--- COURSE MATERIAL ---
{trimmed}
--- END MATERIAL ---

Generate exactly {req.question_count} multiple choice questions based ONLY on the material above.

Rules:
- Each question must have exactly 4 options (A, B, C, D)
- Only one option is correct
- Questions must come directly from the material — no outside knowledge
- Vary the topics covered — don't repeat the same concept twice
- Match the difficulty instruction above

Respond ONLY with a JSON array. No explanation, no markdown, no backticks:
[
  {{
    "id": 1,
    "topic": "Topic this question covers",
    "question": "The question text?",
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "correct_index": 0,
    "explanation": "Brief explanation of why the correct answer is right"
  }}
]

correct_index is 0-based (0 = A, 1 = B, 2 = C, 3 = D).
"""

    raw = await call_gemini(prompt)
    cleaned = raw.strip().replace("```json", "").replace("```", "").strip()

    try:
        questions = json.loads(cleaned)
    except json.JSONDecodeError:
        raise HTTPException(status_code=502, detail=f"Gemini returned invalid JSON: {cleaned[:300]}")

    return {
        "course_code": req.course_code,
        "questions": questions,
        "count": len(questions),
        "difficulty": req.difficulty,
        "material_id": req.material_id,
    }


class ExtractRequest(BaseModel):
    material_id: str
    file_path: str
    file_type: str   # "pdf" | "docx" | "pptx"


@app.post("/extract-text")
async def extract_text(req: ExtractRequest):
    """
    Extracts text from PDF, DOCX, or PPTX files stored in Supabase.
    Called by the frontend after a file is uploaded.
    Replaces the broken Supabase edge functions.
    """
    file_bytes = await fetch_file_from_supabase(req.file_path)

    try:
        if req.file_type == "pdf":
            text = extract_text_from_pdf(file_bytes)
        elif req.file_type == "docx":
            text = extract_text_from_docx(file_bytes)
        elif req.file_type == "pptx":
            text = extract_text_from_pptx(file_bytes)
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {req.file_type}")
    except HTTPException:
        raise
    except Exception as e:
        await save_extracted_content(req.material_id, "", "failed")
        raise HTTPException(status_code=500, detail=f"Extraction failed: {str(e)}")

    if not text or len(text.strip()) < 20:
        await save_extracted_content(req.material_id, "", "scanned_pdf")
        return {"material_id": req.material_id, "status": "scanned_pdf", "chars": 0}

    await save_extracted_content(req.material_id, text)
    return {
        "material_id": req.material_id,
        "status": "success",
        "chars": len(text),
        "preview": text[:200],
    }