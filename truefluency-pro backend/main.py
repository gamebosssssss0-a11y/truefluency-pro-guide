"""
truefluency Backend — FastAPI
Render server that connects the TrueFluency Pro frontend to nvidia nemotron.

Flow:
  Frontend uploads file → Supabase stores it + calls extract edge fn
  Frontend calls /generate-mock → this backend fetches extracted text
  from Supabase → sends to nvidia nemotron → returns real questions to frontend
"""

from fastapi import FastAPI, Header, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from typing import Optional
import os
import json
import io
import httpx
from dotenv import load_dotenv


# File extraction libraries
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

load_dotenv()

app = FastAPI(title="TrueFluency Pro Backend", version="1.0.0")

# ── ENV VARS ──────────────────────────────────────────────────────────────────
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_KEY")   # service role key, NOT the publishable one
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OPENROUTER_MODEL = os.getenv("OPENROUTER_MODEL", "nvidia/nemotron-3-ultra-550b-a55b:free")
PUBLIC_APP_URL = os.getenv("PUBLIC_APP_URL", "https://truefluency-pro.lovable.app")

# Comma-separated browser origins. Empty means "allow all" (dev only).
_origins_raw = (os.getenv("ALLOWED_ORIGINS") or "").strip()
ALLOWED_ORIGINS = [o.strip() for o in _origins_raw.split(",") if o.strip()] or ["*"]

# ── CORS ──────────────────────────────────────────────────────────────────────
# allow_credentials must be False when origins is "*": browsers reject that
# combination outright, and this API is called with a plain fetch (no cookies).
app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=False,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["Content-Type", "Authorization"],
)


def require_env():
    """Fail with a clear message instead of an opaque 502 deep in a handler."""
    missing = [
        name
        for name, value in (
            ("SUPABASE_URL", SUPABASE_URL),
            ("SUPABASE_SERVICE_KEY", SUPABASE_SERVICE_KEY),
            ("OPENROUTER_API_KEY", OPENROUTER_API_KEY),
        )
        if not value
    ]
    if missing:
        raise HTTPException(
            status_code=503,
            detail=f"Server is missing configuration: {', '.join(missing)}",
        )


# ── REQUEST / RESPONSE MODELS ─────────────────────────────────────────────────

class MockRequest(BaseModel):
    material_id: str = Field(min_length=1, max_length=64)
    course_code: str = Field(min_length=1, max_length=32)
    course_name: str = Field(min_length=1, max_length=200)
    question_count: int = Field(default=20, ge=5, le=40)
    difficulty: str = "balanced"  # gentle | balanced | challenging | exam
    topic_focus: list[str] = Field(default_factory=list, max_length=20)

    # User context — built from their onboarding profile
    user_goal: Optional[str] = Field(default=None, max_length=32)
    user_timeline: Optional[str] = Field(default=None, max_length=32)
    user_level: Optional[str] = Field(default=None, max_length=8)
    user_department: Optional[str] = Field(default=None, max_length=120)


class PredictRequest(BaseModel):
    material_id: str = Field(min_length=1, max_length=64)
    course_code: str = Field(min_length=1, max_length=32)
    course_name: str = Field(min_length=1, max_length=200)
    user_level: Optional[str] = Field(default=None, max_length=8)
    user_department: Optional[str] = Field(default=None, max_length=120)



# ── AUTH ──────────────────────────────────────────────────────────────────────

def _service_headers(json_body: bool = True) -> dict:
    headers = {
        "apikey": SUPABASE_SERVICE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_KEY}",
    }
    if json_body:
        headers["Content-Type"] = "application/json"
    return headers


async def require_user(authorization: Optional[str]) -> str:
    """
    Validate the caller's Supabase access token and return their user id.

    Every endpoint that touches a student's private material must go through
    this. Without it, the service role key below would let anyone read or
    overwrite any row in course_materials.
    """
    require_env()
    if not authorization or not authorization.lower().startswith("bearer "):
        raise HTTPException(status_code=401, detail="Sign in to use the analysis service.")

    token = authorization.split(" ", 1)[1].strip()
    if not token:
        raise HTTPException(status_code=401, detail="Sign in to use the analysis service.")

    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            res = await client.get(
                f"{SUPABASE_URL}/auth/v1/user",
                headers={
                    "apikey": SUPABASE_SERVICE_KEY,
                    "Authorization": f"Bearer {token}",
                },
            )
    except httpx.HTTPError:
        raise HTTPException(status_code=503, detail="Could not verify your session. Try again.")

    if res.status_code != 200:
        raise HTTPException(status_code=401, detail="Your session has expired. Sign in again.")

    user_id = (res.json() or {}).get("id")
    if not user_id:
        raise HTTPException(status_code=401, detail="Your session could not be verified.")
    return str(user_id)


# ── HELPERS ───────────────────────────────────────────────────────────────────

async def fetch_owned_material(material_id: str, user_id: str) -> dict:
    """
    Load one course_materials row, scoped to the authenticated owner.

    The user_id filter is what keeps the service role key from becoming an
    IDOR: an id belonging to someone else simply resolves to no rows.
    """
    url = (
        f"{SUPABASE_URL}/rest/v1/course_materials"
        f"?id=eq.{material_id}&user_id=eq.{user_id}"
        "&select=id,user_id,file_path,file_name,file_type,course_code,extracted_content,extraction_status"
    )

    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(url, headers=_service_headers())

    if res.status_code != 200:
        raise HTTPException(status_code=502, detail="Could not load that upload.")

    rows = res.json()
    if not rows:
        # Same response whether the row is missing or owned by someone else,
        # so ids can't be probed for existence.
        raise HTTPException(status_code=404, detail="Material not found")
    return rows[0]


async def fetch_extracted_text(material_id: str, user_id: str) -> str:
    """Extracted text for a material the caller owns."""
    row = await fetch_owned_material(material_id, user_id)

    if row.get("extraction_status") not in ("success",):
        raise HTTPException(
            status_code=422,
            detail=f"Material text not ready yet. Status: {row.get('extraction_status')}"
        )

    content = row.get("extracted_content", "")
    if not content or len(content.strip()) < 50:
        raise HTTPException(status_code=422, detail="Extracted content is too short to generate questions from.")

    return content



async def call_model(prompt: str, max_tokens: int = 4096) -> str:
    """
    Call the configured OpenRouter model and return the raw text response.
    """
    url = "https://openrouter.ai/api/v1/chat/completions"

    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
        "HTTP-Referer": PUBLIC_APP_URL,
        "X-Title": "TrueFluency Pro",
    }

    body = {
        "model": OPENROUTER_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.7,
        "max_tokens": max_tokens,
    }

    try:
        async with httpx.AsyncClient(timeout=120.0) as client:
            res = await client.post(url, headers=headers, json=body)
    except httpx.HTTPError as e:
        raise HTTPException(status_code=504, detail=f"The AI service did not respond in time: {e}")

    if res.status_code != 200:
        raise HTTPException(status_code=502, detail=f"AI service error: {res.text[:300]}")

    data = res.json()
    try:
        return data["choices"][0]["message"]["content"]
    except (KeyError, IndexError, TypeError):
        raise HTTPException(status_code=502, detail="Unexpected AI service response format")


def parse_json_list(raw: str) -> list:
    """
    Models sometimes wrap JSON in markdown or add trailing prose. Strip the
    fences, then fall back to slicing the outermost array before giving up.
    """
    cleaned = raw.strip().replace("```json", "").replace("```", "").strip()
    try:
        parsed = json.loads(cleaned)
    except json.JSONDecodeError:
        start, end = cleaned.find("["), cleaned.rfind("]")
        if start == -1 or end <= start:
            raise HTTPException(
                status_code=502,
                detail=f"The AI returned malformed data: {cleaned[:200]}",
            )
        try:
            parsed = json.loads(cleaned[start : end + 1])
        except json.JSONDecodeError:
            raise HTTPException(
                status_code=502,
                detail=f"The AI returned malformed data: {cleaned[:200]}",
            )
    if not isinstance(parsed, list):
        raise HTTPException(status_code=502, detail="The AI returned an unexpected shape.")
    return parsed



def build_user_context(req) -> str:
    """
    Build a user context string from the profile data sent by the frontend.
    This is what makes the AI understand WHO it's generating for.
    """
    goal_map = {
        "pass": "just needs to pass (focus on high-frequency topics only)",
        "top-grades": "aiming for top grades (cover everything including edge topics)",
        "catch-up": "catching up on missed content (cover fundamentals thoroughly)",
    }
    timeline_map = {
        "lt-week": "less than a week until their exam (urgent — prioritize the most likely topics)",
        "2-4-weeks": "2 to 4 weeks until their exam",
        "gt-month": "more than a month until their exam (can afford more breadth)",
        "unsure": "unsure of their exam date",
    }

    parts = []
    if req.user_level:
        parts.append(f"- Academic level: {req.user_level}L at University of Ibadan")
    if req.user_department:
        parts.append(f"- Department: {req.user_department}")
    if req.user_goal:
        parts.append(f"- Goal: Student {goal_map.get(req.user_goal, req.user_goal)}")
    if req.user_timeline:
        parts.append(f"- Timeline: {timeline_map.get(req.user_timeline, req.user_timeline)}")

    if not parts:
        return ""

    return "Student context:\n" + "\n".join(parts)


def difficulty_instruction(difficulty: str) -> str:
    return {
        "gentle": "Make the questions straightforward and accessible. Focus on definitions, basic concepts, and direct recall.",
        "balanced": "Mix easy recall questions with questions that require understanding and application.",
        "challenging": "Include questions that require deeper understanding, analysis, and application of concepts.",
        "exam": "Match real exam difficulty exactly. Include tricky options, edge cases, and questions requiring synthesis.",
    }.get(difficulty, "Mix easy and hard questions.")


# ── ROUTES ────────────────────────────────────────────────────────────────────

async def fetch_file_from_supabase(file_path: str) -> bytes:
    """Download a file from Supabase storage and return raw bytes."""
    url = f"{SUPABASE_URL}/storage/v1/object/course-materials/{file_path}"
    headers = {
        "apikey": SUPABASE_SERVICE_KEY,
        "Authorization": f"Bearer {SUPABASE_SERVICE_KEY}",
    }
    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.get(url, headers=headers)
    if res.status_code != 200:
        raise HTTPException(status_code=502, detail=f"Could not fetch file from storage: {res.text}")
    return res.content


def extract_text_from_pdf(file_bytes: bytes) -> str:
    if not HAS_PDF:
        raise HTTPException(status_code=500, detail="PyPDF2 not installed")
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        text += (page.extract_text() or "") + "\n"
    return text.strip()


def extract_text_from_docx(file_bytes: bytes) -> str:
    if not HAS_DOCX:
        raise HTTPException(status_code=500, detail="python-docx not installed")
    doc = Document(io.BytesIO(file_bytes))
    text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())
    return text.strip()


def extract_text_from_pptx(file_bytes: bytes) -> str:
    if not HAS_PPTX:
        raise HTTPException(status_code=500, detail="python-pptx not installed")
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text.strip()


async def save_extracted_content(
    material_id: str, user_id: str, content: str, status: str = "success"
):
    """Save extracted text back to the caller's own course_materials row."""
    url = (
        f"{SUPABASE_URL}/rest/v1/course_materials"
        f"?id=eq.{material_id}&user_id=eq.{user_id}"
    )
    body = {"extracted_content": content, "extraction_status": status}
    async with httpx.AsyncClient(timeout=30.0) as client:
        res = await client.patch(url, headers=_service_headers(), json=body)
    if res.status_code not in (200, 204):
        raise HTTPException(status_code=502, detail="Failed to save extraction.")


@app.get("/health")
def health():
    """Quick check that the server is alive — call this first when testing."""
    return {
        "status": "ok",
        "service": "TrueFluency Pro Backend",
        "configured": bool(SUPABASE_URL and SUPABASE_SERVICE_KEY and OPENROUTER_API_KEY),
    }


@app.post("/predict-topics")
async def predict_topics(
    req: PredictRequest,
    authorization: Optional[str] = Header(default=None),
):
    """
    Step 1 of the AI pipeline.
    Takes a material_id, fetches the extracted text from Supabase,
    sends it to the model, returns predicted exam topics with confidence scores.

    Frontend uses this to replace the hardcoded placeholder topics on the
    dashboard and course detail screens.
    """
    user_id = await require_user(authorization)
    extracted_text = await fetch_extracted_text(req.material_id, user_id)

    # Trim to avoid token overload — first 6000 chars is usually enough for topic prediction
    trimmed = extracted_text[:6000]


    prompt = f"""
You are an exam prediction engine for University of Ibadan ({req.user_level or "undergraduate"} level).

Analyze this course material for {req.course_code} — {req.course_name}:
{req.user_department and f"Department: {req.user_department}" or ""}

--- MATERIAL START ---
{trimmed}
--- MATERIAL END ---

Based on this content, identify the most likely exam topics.
For each topic, give a confidence score (0.0 to 1.0) based on how prominently it features in the material.

Respond ONLY with a JSON array. No explanation, no markdown, no backticks. Just raw JSON like:
[
  {{"topic": "Topic Name", "confidence": 0.9}},
  {{"topic": "Another Topic", "confidence": 0.6}}
]

Return 5 to 8 topics, ordered from highest to lowest confidence.
"""

    raw = await call_model(prompt, max_tokens=1500)
    topics = parse_json_list(raw)


    return {
        "course_code": req.course_code,
        "topics": topics,
        "material_id": req.material_id,
    }


@app.post("/generate-mock")
async def generate_mock(
    req: MockRequest,
    authorization: Optional[str] = Header(default=None),
):
    """
    Step 2 of the AI pipeline — the core endpoint.
    Generates a full mock test based on the extracted material.

    This replaces the hardcoded sampleQuestions from questions.ts in the frontend.
    """
    user_id = await require_user(authorization)
    extracted_text = await fetch_extracted_text(req.material_id, user_id)


    # Trim to keep within Gemini's context window comfortably
    trimmed = extracted_text[:8000]

    user_ctx = build_user_context(req)
    diff_instruction = difficulty_instruction(req.difficulty)

    topic_instruction = ""
    if req.topic_focus:
        topic_instruction = f"Pay extra attention to these topics (weight them more heavily): {', '.join(req.topic_focus)}."

    prompt = f"""
You are a mock exam generator for University of Ibadan.

{user_ctx}

Course: {req.course_code} — {req.course_name}
Difficulty: {req.difficulty}. {diff_instruction}
{topic_instruction}

--- COURSE MATERIAL ---
{trimmed}
--- END MATERIAL ---

Generate exactly {req.question_count} multiple choice questions based ONLY on the material above.

Rules:
- Each question must have exactly 4 options (A, B, C, D)
- Only one option is correct
- Questions must come directly from the material — no outside knowledge
- Vary the topics covered — don't repeat the same concept twice
- Match the difficulty instruction above

Respond ONLY with a JSON array. No explanation, no markdown, no backticks:
[
  {{
    "id": 1,
    "topic": "Topic this question covers",
    "question": "The question text?",
    "options": ["Option A", "Option B", "Option C", "Option D"],
    "correct_index": 0,
    "explanation": "Brief explanation of why the correct answer is right"
  }}
]

correct_index is 0-based (0 = A, 1 = B, 2 = C, 3 = D).
"""

    # ~350 tokens per MCQ with options + explanation, plus headroom, so a
    # 40-question set can't silently truncate into malformed JSON.
    raw = await call_model(prompt, max_tokens=min(16000, 600 + req.question_count * 400))
    questions = parse_json_list(raw)


    return {
        "course_code": req.course_code,
        "questions": questions,
        "count": len(questions),
        "difficulty": req.difficulty,
        "material_id": req.material_id,
    }


class ExtractRequest(BaseModel):
    material_id: str = Field(min_length=1, max_length=64)
    file_type: str = Field(min_length=1, max_length=8)   # "pdf" | "docx" | "pptx"


@app.post("/extract-text")
async def extract_text(
    req: ExtractRequest,
    authorization: Optional[str] = Header(default=None),
):
    """
    Extracts text from PDF, DOCX, or PPTX files stored in Supabase.
    Called by the frontend after a file is uploaded.

    The storage path is read from the caller's own material row, never taken
    from the request body, so an arbitrary path can't be pulled out of the
    private bucket.
    """
    user_id = await require_user(authorization)
    row = await fetch_owned_material(req.material_id, user_id)

    file_path = row.get("file_path") or ""
    if not file_path or not file_path.startswith(f"{user_id}/"):
        raise HTTPException(status_code=404, detail="Material not found")

    file_bytes = await fetch_file_from_supabase(file_path)

    try:
        if req.file_type == "pdf":
            text = extract_text_from_pdf(file_bytes)
        elif req.file_type == "docx":
            text = extract_text_from_docx(file_bytes)
        elif req.file_type == "pptx":
            text = extract_text_from_pptx(file_bytes)
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {req.file_type}")
    except HTTPException:
        raise
    except Exception as e:
        await save_extracted_content(req.material_id, user_id, "", "failed")
        print(f"[extract-text] extraction failed for {req.material_id}: {e}")
        raise HTTPException(status_code=500, detail="We couldn't read that file.")

    if not text or len(text.strip()) < 20:
        await save_extracted_content(req.material_id, user_id, "", "scanned_pdf")
        return {"material_id": req.material_id, "status": "scanned_pdf", "chars": 0}

    await save_extracted_content(req.material_id, user_id, text)
    return {
        "material_id": req.material_id,
        "status": "success",
        "chars": len(text),
        "preview": text[:200],
    }

